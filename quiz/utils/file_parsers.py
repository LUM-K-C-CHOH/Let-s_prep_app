import os
import tempfile

from PyPDF2 import PdfReader
import docx  # from python-docx
from pptx import Presentation


def extract_text_from_file(django_file):
    """
    Accepts a Django UploadedFile, saves it to a temp file,
    and extracts text depending on extension.
    """
    name = django_file.name.lower()
    ext = os.path.splitext(name)[1]

    # Save upload to a temp file
    with tempfile.NamedTemporaryFile(delete=False) as tmp:
        for chunk in django_file.chunks():
            tmp.write(chunk)
        temp_path = tmp.name

    try:
        if ext == '.pdf':
            return _extract_pdf(temp_path)
        elif ext in ['.docx']:
            return _extract_docx(temp_path)
        elif ext in ['.pptx']:
            return _extract_pptx(temp_path)
        elif ext in ['.txt']:
            return _extract_txt(temp_path)
        else:
            # Fallback: try reading as text
            return _extract_txt(temp_path)
    finally:
        try:
            os.remove(temp_path)
        except FileNotFoundError:
            pass


def _extract_pdf(path):
    text = ""
    reader = PdfReader(path)
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def _extract_docx(path):
    doc = docx.Document(path)
    text = "\n".join(p.text for p in doc.paragraphs)
    return text


def _extract_pptx(path):
    prs = Presentation(path)
    text_chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_chunks.append(shape.text)
    return "\n".join(text_chunks)


def _extract_txt(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(path, "r", encoding="latin-1", errors="ignore") as f:
            return f.read()